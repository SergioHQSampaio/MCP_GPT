from __future__ import annotations

import base64
import io
import json
import mimetypes
import os
import posixpath
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from dotenv import load_dotenv
from google.oauth2 import service_account
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaInMemoryUpload
from mcp.server.fastmcp import FastMCP

load_dotenv()

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/documents",
    "https://www.googleapis.com/auth/spreadsheets",
]

MCP_HOST = os.environ.get("MCP_HOST", "0.0.0.0")
MCP_PORT = int(os.environ.get("MCP_PORT", "8787"))
MCP_MOUNT_PATH = os.environ.get("MCP_MOUNT_PATH", "/")
MCP_STREAMABLE_HTTP_PATH = os.environ.get("MCP_STREAMABLE_HTTP_PATH", "/mcp")
ROOT_FOLDER_ID = os.environ.get("GDRIVE_ROOT_FOLDER_ID", "").strip()
DEFAULT_TEXT_MIME = "text/markdown"

if not ROOT_FOLDER_ID:
    raise RuntimeError("GDRIVE_ROOT_FOLDER_ID nao configurado no ambiente.")

SERVER_INSTRUCTIONS = """
Este conector opera exclusivamente dentro de uma pasta-raiz especifica do Google Drive.

Trate-o como uma camada geral de operacao de arquivos na subarvore permitida.
Nao trate Google Docs, Google Sheets ou extensoes especificas como a definicao do produto.

Capacidades principais:
- listar, localizar, ler metadados, criar, substituir, mover, renomear, enviar para a lixeira e restaurar arquivos e pastas;
- ler semanticamente formatos suportados, escolhendo o backend mais adequado por MIME ou extensao;
- preservar o formato pedido pelo usuario quando isso for tecnicamente viavel;
- usar formatos Google nativos apenas quando forem o melhor meio tecnico ou quando o usuario pedir explicitamente;
- diferenciar leitura semantica, escrita semantica e simples substituicao binaria.

Regras obrigatorias:
- Nunca opere fora da pasta-raiz configurada.
- Nunca use IDs externos sem validar se pertencem a subarvore permitida.
- Quando o usuario informar apenas um caminho relativo, resolva sempre a partir da pasta-raiz.
- Quando a edicao interna do formato for fraca ou insegura, use reconstrucao ou substituicao de arquivo.
- Para o mesmo pedido, prefira o nivel maximo de operacao tecnicamente viavel para o formato solicitado.
""".strip()

mcp = FastMCP(
    name="Google Drive Bridge",
    instructions=SERVER_INSTRUCTIONS,
    host=MCP_HOST,
    port=MCP_PORT,
    mount_path=MCP_MOUNT_PATH,
    streamable_http_path=MCP_STREAMABLE_HTTP_PATH,
)


try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from openpyxl import Workbook, load_workbook
except Exception:
    Workbook = None
    load_workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from PIL import Image
except Exception:
    Image = None


GOOGLE_DOC_MIME = "application/vnd.google-apps.document"
GOOGLE_SHEET_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDES_MIME = "application/vnd.google-apps.presentation"
FOLDER_MIME = "application/vnd.google-apps.folder"

TEXT_EXTENSIONS = {".txt", ".md", ".json", ".csv", ".html", ".xml", ".yaml", ".yml"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".svg"}

EXPORT_MIME_BY_HINT = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "txt": "text/plain",
    "md": "text/markdown",
    "csv": "text/csv",
}

GOOGLE_NATIVE_TARGETS = {
    "gdoc": GOOGLE_DOC_MIME,
    "google_doc": GOOGLE_DOC_MIME,
    "google-doc": GOOGLE_DOC_MIME,
    "gsheet": GOOGLE_SHEET_MIME,
    "google_sheet": GOOGLE_SHEET_MIME,
    "google-sheet": GOOGLE_SHEET_MIME,
    "gslides": GOOGLE_SLIDES_MIME,
    "google_slides": GOOGLE_SLIDES_MIME,
    "google-slides": GOOGLE_SLIDES_MIME,
}


@dataclass
class GoogleClients:
    drive: Any
    docs: Any
    sheets: Any
    slides: Any


def _slug(name: str) -> str:
    name = (name or "").strip()
    name = re.sub(r"[\\/:*?\"<>|]+", "-", name)
    name = re.sub(r"\s+", " ", name).strip()
    return name or "item"


def _ext(name: str) -> str:
    return Path(name or "").suffix.lower()


def _build_credentials() -> Credentials:
    service_account_json = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON", "").strip()
    service_account_file = os.environ.get("GOOGLE_SERVICE_ACCOUNT_FILE", "").strip()
    delegated_user = os.environ.get("GOOGLE_DELEGATED_USER", "").strip()

    if service_account_json:
        info = json.loads(service_account_json)
        creds = service_account.Credentials.from_service_account_info(info, scopes=SCOPES)
        if delegated_user:
            creds = creds.with_subject(delegated_user)
        return creds

    if service_account_file:
        creds = service_account.Credentials.from_service_account_file(service_account_file, scopes=SCOPES)
        if delegated_user:
            creds = creds.with_subject(delegated_user)
        return creds

    refresh_token = os.environ.get("GOOGLE_REFRESH_TOKEN", "").strip()
    client_id = os.environ.get("GOOGLE_CLIENT_ID", "").strip()
    client_secret = os.environ.get("GOOGLE_CLIENT_SECRET", "").strip()
    token_uri = os.environ.get("GOOGLE_TOKEN_URI", "https://oauth2.googleapis.com/token").strip()

    if refresh_token and client_id and client_secret:
        creds = Credentials(
            token=None,
            refresh_token=refresh_token,
            token_uri=token_uri,
            client_id=client_id,
            client_secret=client_secret,
            scopes=SCOPES,
        )
        creds.refresh(Request())
        return creds

    raise RuntimeError(
        "Configure autenticacao Google: GOOGLE_SERVICE_ACCOUNT_FILE / GOOGLE_SERVICE_ACCOUNT_JSON "
        "ou GOOGLE_REFRESH_TOKEN + GOOGLE_CLIENT_ID + GOOGLE_CLIENT_SECRET."
    )


def _clients() -> GoogleClients:
    creds = _build_credentials()
    return GoogleClients(
        drive=build("drive", "v3", credentials=creds, cache_discovery=False),
        docs=build("docs", "v1", credentials=creds, cache_discovery=False),
        sheets=build("sheets", "v4", credentials=creds, cache_discovery=False),
        slides=build("slides", "v1", credentials=creds, cache_discovery=False),
    )


def _root_metadata(clients: GoogleClients) -> dict[str, Any]:
    return clients.drive.files().get(
        fileId=ROOT_FOLDER_ID,
        fields="id,name,mimeType,parents,trashed,webViewLink",
        supportsAllDrives=True,
    ).execute()


def _normalize_relpath(path: str | None) -> str:
    raw = (path or "").strip().replace("\\", "/")
    raw = posixpath.normpath(raw)
    if raw in (".", "/"):
        return ""
    if raw.startswith("../") or raw == "..":
        raise ValueError("Caminho relativo invalido; sai da pasta-raiz.")
    return raw.strip("/")


def _get_file(clients: GoogleClients, file_id: str, fields: str = "id,name,mimeType,parents,trashed") -> dict[str, Any]:
    return clients.drive.files().get(
        fileId=file_id,
        fields=fields,
        supportsAllDrives=True,
    ).execute()


def _list_children(clients: GoogleClients, parent_id: str) -> list[dict[str, Any]]:
    results = []
    page_token = None
    while True:
        resp = clients.drive.files().list(
            q=f"'{parent_id}' in parents and trashed = false",
            fields="nextPageToken, files(id,name,mimeType,parents,modifiedTime,size,webViewLink,iconLink)",
            pageToken=page_token,
            pageSize=1000,
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
        ).execute()
        results.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return results


def _child_by_name(clients: GoogleClients, parent_id: str, name: str) -> dict[str, Any] | None:
    escaped = name.replace("'", "\\'")
    resp = clients.drive.files().list(
        q=f"'{parent_id}' in parents and name = '{escaped}' and trashed = false",
        fields="files(id,name,mimeType,parents,modifiedTime,size,webViewLink)",
        pageSize=10,
        includeItemsFromAllDrives=True,
        supportsAllDrives=True,
    ).execute()
    files = resp.get("files", [])
    return files[0] if files else None


def _resolve_path(clients: GoogleClients, path: str | None, create_missing_folders: bool = False, final_folder: bool = True) -> dict[str, Any]:
    rel = _normalize_relpath(path)
    current = _root_metadata(clients)
    if not rel:
        return current
    parts = [p for p in rel.split("/") if p]
    for idx, part in enumerate(parts):
        child = _child_by_name(clients, current["id"], part)
        is_last = idx == len(parts) - 1
        if child is None:
            if create_missing_folders or (final_folder and is_last):
                child = clients.drive.files().create(
                    body={"name": part, "mimeType": FOLDER_MIME, "parents": [current["id"]]},
                    fields="id,name,mimeType,parents,webViewLink",
                    supportsAllDrives=True,
                ).execute()
            else:
                raise FileNotFoundError(f"Caminho nao encontrado: {rel}")
        current = child
    return current


def _ensure_inside_root(clients: GoogleClients, file_id: str) -> dict[str, Any]:
    current = _get_file(clients, file_id, fields="id,name,mimeType,parents,trashed,webViewLink,modifiedTime,size")
    visited = set()
    probe = current
    while True:
        if probe["id"] == ROOT_FOLDER_ID:
            return current
        if probe["id"] in visited:
            break
        visited.add(probe["id"])
        parents = probe.get("parents") or []
        if not parents:
            break
        probe = _get_file(clients, parents[0], fields="id,name,mimeType,parents")
    raise PermissionError("O item solicitado esta fora da pasta-raiz permitida.")


def _metadata(file_obj: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": file_obj.get("id"),
        "name": file_obj.get("name"),
        "mimeType": file_obj.get("mimeType"),
        "parents": file_obj.get("parents"),
        "modifiedTime": file_obj.get("modifiedTime"),
        "size": file_obj.get("size"),
        "webViewLink": file_obj.get("webViewLink"),
    }


def _doc_text(clients: GoogleClients, document_id: str) -> str:
    doc = clients.docs.documents().get(documentId=document_id).execute()
    chunks: list[str] = []
    for el in doc.get("body", {}).get("content", []):
        para = el.get("paragraph")
        if not para:
            continue
        for pe in para.get("elements", []):
            text_run = pe.get("textRun")
            if text_run and text_run.get("content"):
                chunks.append(text_run["content"])
    return "".join(chunks)


def _clear_and_write_doc(clients: GoogleClients, document_id: str, content: str) -> None:
    doc = clients.docs.documents().get(documentId=document_id).execute()
    end_index = doc.get("body", {}).get("content", [{}])[-1].get("endIndex", 1)
    requests = []
    if end_index and end_index > 2:
        requests.append({"deleteContentRange": {"range": {"startIndex": 1, "endIndex": end_index - 1}}})
    if content:
        requests.append({"insertText": {"location": {"index": 1}, "text": content}})
    if requests:
        clients.docs.documents().batchUpdate(documentId=document_id, body={"requests": requests}).execute()


def _sheet_values(clients: GoogleClients, spreadsheet_id: str, a1_range: str) -> list[list[Any]]:
    resp = clients.sheets.spreadsheets().values().get(spreadsheetId=spreadsheet_id, range=a1_range).execute()
    return resp.get("values", [])


def _sheet_write(clients: GoogleClients, spreadsheet_id: str, a1_range: str, values: list[list[Any]]) -> None:
    clients.sheets.spreadsheets().values().update(
        spreadsheetId=spreadsheet_id,
        range=a1_range,
        valueInputOption="USER_ENTERED",
        body={"values": values},
    ).execute()


def _guess_mime(filename: str, fallback: str = "application/octet-stream") -> str:
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or fallback


def _download_bytes(clients: GoogleClients, file_id: str) -> bytes:
    request = clients.drive.files().get_media(fileId=file_id, supportsAllDrives=True)
    buffer = io.BytesIO()
    downloader = __import__("googleapiclient.http").http.MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buffer.getvalue()


def _upload_or_replace_bytes(clients: GoogleClients, parent: dict[str, Any], filename: str, content: bytes, mime_type: str, replace_if_exists: bool = True) -> dict[str, Any]:
    filename = _slug(filename)
    media = MediaInMemoryUpload(content, mimetype=mime_type, resumable=False)
    existing = _child_by_name(clients, parent["id"], filename)
    if existing:
        if not replace_if_exists:
            raise ValueError("O arquivo ja existe e replace_if_exists=false.")
        updated = clients.drive.files().update(
            fileId=existing["id"],
            media_body=media,
            fields="id,name,mimeType,parents,webViewLink,modifiedTime,size",
            supportsAllDrives=True,
        ).execute()
        return {"ok": True, "created": False, "file": _metadata(updated)}
    created = clients.drive.files().create(
        body={"name": filename, "parents": [parent["id"]]},
        media_body=media,
        fields="id,name,mimeType,parents,webViewLink,modifiedTime,size",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "created": True, "file": _metadata(created)}


def _require_library(lib: Any, label: str) -> None:
    if lib is None:
        raise RuntimeError(f"Biblioteca ausente para {label}. Adicione a dependencia correspondente no ambiente.")


def _serialize_docx(content: Any) -> bytes:
    _require_library(DocxDocument, "DOCX")
    doc = DocxDocument()
    if isinstance(content, str):
        for block in content.splitlines() or [""]:
            doc.add_paragraph(block)
    elif isinstance(content, list):
        for item in content:
            doc.add_paragraph(str(item))
    elif isinstance(content, dict):
        for item in content.get("paragraphs", []):
            doc.add_paragraph(str(item))
    else:
        doc.add_paragraph(str(content))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _parse_docx(binary: bytes) -> dict[str, Any]:
    _require_library(DocxDocument, "DOCX")
    doc = DocxDocument(io.BytesIO(binary))
    paragraphs = [p.text for p in doc.paragraphs]
    return {"kind": "docx", "paragraphs": paragraphs, "text": "\n".join(paragraphs)}


def _serialize_xlsx(content: Any) -> bytes:
    _require_library(Workbook, "XLSX")
    wb = Workbook()
    default_sheet = wb.active
    if isinstance(content, dict) and content.get("sheets"):
        first = True
        for sheet_name, rows in content["sheets"].items():
            ws = default_sheet if first else wb.create_sheet()
            first = False
            ws.title = _slug(str(sheet_name))[:31]
            for row in rows:
                ws.append(list(row))
    elif isinstance(content, list):
        default_sheet.title = "Sheet1"
        for row in content:
            default_sheet.append(list(row) if isinstance(row, (list, tuple)) else [row])
    else:
        default_sheet["A1"] = str(content)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _parse_xlsx(binary: bytes) -> dict[str, Any]:
    _require_library(load_workbook, "XLSX")
    wb = load_workbook(io.BytesIO(binary), data_only=True)
    sheets: dict[str, list[list[Any]]] = {}
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))
        sheets[ws.title] = rows
    return {"kind": "xlsx", "sheets": sheets}


def _serialize_pptx(content: Any) -> bytes:
    _require_library(Presentation, "PPTX")
    prs = Presentation()
    slides = content.get("slides") if isinstance(content, dict) else content
    if not isinstance(slides, list):
        slides = [{"title": "Slide 1", "body": str(content)}]
    for slide_data in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title = slide_data.get("title", "") if isinstance(slide_data, dict) else str(slide_data)
        body = slide_data.get("body", "") if isinstance(slide_data, dict) else ""
        slide.shapes.title.text = str(title)
        slide.placeholders[1].text = str(body)
    if len(prs.slides) > 0 and len(content if isinstance(content, list) else []) == 0:
        pass
    if len(prs.slides) > 1 and prs.slides[0].shapes.title.text == "":
        pass
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _parse_pptx(binary: bytes) -> dict[str, Any]:
    _require_library(Presentation, "PPTX")
    prs = Presentation(io.BytesIO(binary))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slides.append({"index": idx, "text": texts})
    return {"kind": "pptx", "slides": slides}


def _parse_pdf(binary: bytes) -> dict[str, Any]:
    _require_library(PdfReader, "PDF")
    reader = PdfReader(io.BytesIO(binary))
    pages = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({"page": idx, "text": text})
    return {"kind": "pdf", "pages": pages, "text": "\n\n".join(p["text"] for p in pages)}


def _parse_image(binary: bytes, name: str, mime_type: str) -> dict[str, Any]:
    info = {"kind": "image", "filename": name, "mimeType": mime_type, "bytes": len(binary)}
    if Image is not None:
        try:
            img = Image.open(io.BytesIO(binary))
            info.update({
                "format": img.format,
                "mode": img.mode,
                "width": img.width,
                "height": img.height,
            })
        except Exception:
            pass
    return info


def _detect_handler(file_obj: dict[str, Any]) -> str:
    mime_type = (file_obj.get("mimeType") or "").lower()
    ext = _ext(file_obj.get("name", ""))
    if mime_type == GOOGLE_DOC_MIME:
        return "google_doc"
    if mime_type == GOOGLE_SHEET_MIME:
        return "google_sheet"
    if mime_type == GOOGLE_SLIDES_MIME:
        return "google_slides"
    if ext in TEXT_EXTENSIONS or mime_type.startswith("text/") or mime_type in {"application/json", "application/xml"}:
        return "plain_text"
    if ext == ".docx":
        return "docx"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".pdf" or mime_type == "application/pdf":
        return "pdf"
    if ext in IMAGE_EXTENSIONS or mime_type.startswith("image/"):
        return "image"
    return "binary"


def _read_google_slides_summary(clients: GoogleClients, file_id: str) -> dict[str, Any]:
    presentation = clients.slides.presentations().get(presentationId=file_id).execute()
    slides_out = []
    for idx, slide in enumerate(presentation.get("slides", []), start=1):
        texts: list[str] = []
        notes: list[str] = []
        for element in slide.get("pageElements", []):
            shape = element.get("shape") or {}
            text_content = shape.get("text") or {}
            for te in text_content.get("textElements", []):
                run = te.get("textRun") or {}
                content = (run.get("content") or "").strip()
                if content:
                    texts.append(content)
        notes_page = slide.get("slideProperties", {}).get("notesPage") or {}
        for element in notes_page.get("pageElements", []):
            shape = element.get("shape") or {}
            text_content = shape.get("text") or {}
            for te in text_content.get("textElements", []):
                run = te.get("textRun") or {}
                content = (run.get("content") or "").strip()
                if content:
                    notes.append(content)
        slides_out.append({"index": idx, "objectId": slide.get("objectId"), "text": texts, "notes": notes})
    return {"kind": "google_slides", "title": presentation.get("title"), "slides": slides_out}


def _write_google_slides(clients: GoogleClients, presentation_id: str, slides: Any) -> None:
    if not isinstance(slides, list):
        if isinstance(slides, dict) and "slides" in slides:
            slides = slides.get("slides")
        else:
            slides = [{"title": "Slide 1", "body": str(slides)}]
    if not slides:
        slides = [{"title": "Slide 1", "body": ""}]

    presentation = clients.slides.presentations().get(presentationId=presentation_id).execute()
    existing_ids = [s.get("objectId") for s in presentation.get("slides", []) if s.get("objectId")]
    requests = []
    if existing_ids:
        requests.extend({"deleteObject": {"objectId": oid}} for oid in existing_ids)

    for idx, slide_data in enumerate(slides, start=1):
        title = slide_data.get("title", f"Slide {idx}") if isinstance(slide_data, dict) else f"Slide {idx}"
        body = slide_data.get("body", "") if isinstance(slide_data, dict) else str(slide_data)
        sid = f"slide_{idx}"
        tid = f"title_{idx}"
        bid = f"body_{idx}"
        requests.append({
            "createSlide": {
                "objectId": sid,
                "slideLayoutReference": {"predefinedLayout": "TITLE_AND_BODY"},
                "placeholderIdMappings": [
                    {"layoutPlaceholder": {"type": "TITLE", "index": 0}, "objectId": tid},
                    {"layoutPlaceholder": {"type": "BODY", "index": 0}, "objectId": bid},
                ],
            }
        })
        if title:
            requests.append({"insertText": {"objectId": tid, "insertionIndex": 0, "text": str(title)}})
        if body:
            requests.append({"insertText": {"objectId": bid, "insertionIndex": 0, "text": str(body)}})

    if requests:
        clients.slides.presentations().batchUpdate(presentationId=presentation_id, body={"requests": requests}).execute()


def _read_google_sheet_summary(clients: GoogleClients, file_id: str) -> dict[str, Any]:
    spreadsheet = clients.sheets.spreadsheets().get(spreadsheetId=file_id).execute()
    sheet_titles = [s.get("properties", {}).get("title") for s in spreadsheet.get("sheets", [])]
    preview: dict[str, list[list[Any]]] = {}
    for title in sheet_titles[:3]:
        if title:
            preview[title] = _sheet_values(clients, file_id, f"'{title}'!A1:Z20")
    return {"kind": "google_sheet", "sheetTitles": sheet_titles, "preview": preview}


def _read_file_auto(clients: GoogleClients, file_obj: dict[str, Any]) -> dict[str, Any]:
    handler = _detect_handler(file_obj)
    if handler == "google_doc":
        return {"kind": handler, "text": _doc_text(clients, file_obj["id"])}
    if handler == "google_sheet":
        return _read_google_sheet_summary(clients, file_obj["id"])
    if handler == "google_slides":
        return _read_google_slides_summary(clients, file_obj["id"])

    binary = _download_bytes(clients, file_obj["id"])
    if handler == "plain_text":
        return {"kind": handler, "text": binary.decode("utf-8", errors="replace")}
    if handler == "docx":
        return _parse_docx(binary)
    if handler == "xlsx":
        return _parse_xlsx(binary)
    if handler == "pptx":
        return _parse_pptx(binary)
    if handler == "pdf":
        return _parse_pdf(binary)
    if handler == "image":
        return _parse_image(binary, file_obj.get("name", ""), file_obj.get("mimeType", ""))
    return {"kind": "binary", "message": "Formato sem leitura semantica implementada nesta fase.", "bytes": len(binary)}


def _serialize_content_for_filename(filename: str, content: Any, mime_type: str = "") -> tuple[bytes, str]:
    ext = _ext(filename)
    guessed = mime_type or _guess_mime(filename, DEFAULT_TEXT_MIME)

    if ext in {".txt", ".md", ".json", ".csv", ".html", ".xml", ".yaml", ".yml"} or guessed.startswith("text/") or guessed in {"application/json", "application/xml"}:
        if isinstance(content, (dict, list)) and ext == ".json":
            return json.dumps(content, ensure_ascii=False, indent=2).encode("utf-8"), guessed
        return str(content).encode("utf-8"), guessed
    if ext == ".docx":
        return _serialize_docx(content), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if ext == ".xlsx":
        return _serialize_xlsx(content), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if ext == ".pptx":
        return _serialize_pptx(content), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if isinstance(content, dict) and "base64_content" in content:
        return base64.b64decode(content["base64_content"]), mime_type or _guess_mime(filename)
    if isinstance(content, str) and re.fullmatch(r"[A-Za-z0-9+/=\s]+", content or "") and guessed == "application/octet-stream":
        try:
            return base64.b64decode(content), guessed
        except Exception:
            pass
    return str(content).encode("utf-8"), guessed


def _create_google_slides(clients: GoogleClients, parent_path: str, title: str, slides: Any = None) -> dict[str, Any]:
    parent = _resolve_path(clients, parent_path, final_folder=False)
    created = clients.drive.files().create(
        body={"name": _slug(title), "mimeType": GOOGLE_SLIDES_MIME, "parents": [parent["id"]]},
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    if slides is not None:
        _write_google_slides(clients, created["id"], slides)
    summary = _read_google_slides_summary(clients, created["id"])
    return {"ok": True, "presentation": _metadata(created), "slidesRequested": slides if slides is not None else [], "summary": summary}


def _export_google_native_file(clients: GoogleClients, file_obj: dict[str, Any], export_mime: str) -> bytes:
    request = clients.drive.files().export_media(fileId=file_obj["id"], mimeType=export_mime)
    buffer = io.BytesIO()
    downloader = __import__("googleapiclient.http").http.MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buffer.getvalue()


def _pick_export_mime(file_obj: dict[str, Any], target_format: str = "", mime_type: str = "") -> str:
    target = (target_format or "").lower().strip()
    if mime_type:
        return mime_type
    if target in EXPORT_MIME_BY_HINT:
        return EXPORT_MIME_BY_HINT[target]
    source_mime = (file_obj.get("mimeType") or "").lower()
    if source_mime == GOOGLE_DOC_MIME:
        return EXPORT_MIME_BY_HINT["docx"]
    if source_mime == GOOGLE_SHEET_MIME:
        return EXPORT_MIME_BY_HINT["xlsx"]
    if source_mime == GOOGLE_SLIDES_MIME:
        return EXPORT_MIME_BY_HINT["pptx"]
    raise ValueError("Nao foi possivel determinar o mime_type de exportacao. Informe target_format ou mime_type.")


def _guess_filename_for_conversion(file_obj: dict[str, Any], target_format: str = "", mime_type: str = "") -> str:
    source_name = file_obj.get("name") or "item"
    stem = Path(source_name).stem or source_name
    target = (target_format or "").lower().strip()
    if target in {"gdoc", "google_doc", "google-doc", "gsheet", "google_sheet", "google-sheet", "gslides", "google_slides", "google-slides"}:
        return stem
    if target:
        return f"{stem}.{target}"
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return f"{stem}.docx"
    if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return f"{stem}.xlsx"
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return f"{stem}.pptx"
    if mime_type == "application/pdf":
        return f"{stem}.pdf"
    if mime_type == "text/plain":
        return f"{stem}.txt"
    if mime_type == "text/markdown":
        return f"{stem}.md"
    if mime_type == "text/csv":
        return f"{stem}.csv"
    return stem


def _materialize_converted_content(clients: GoogleClients, file_obj: dict[str, Any], target_format: str = "", mime_type: str = "") -> tuple[bytes, str, str]:
    source_mime = (file_obj.get("mimeType") or "").lower()
    target = (target_format or "").lower().strip()

    if target in GOOGLE_NATIVE_TARGETS:
        logical = _read_file_auto(clients, file_obj)
        if target in {"gdoc", "google_doc", "google-doc"}:
            if logical.get("kind") in {"google_doc", "plain_text", "docx", "pdf"}:
                text = logical.get("text") or ""
            elif logical.get("kind") == "pptx":
                text = "\n\n".join("\n".join(slide.get("text", [])) for slide in logical.get("slides", []))
            else:
                raise ValueError("Conversao para Google Doc nao implementada para este formato.")
            return text.encode("utf-8"), GOOGLE_NATIVE_TARGETS[target], _guess_filename_for_conversion(file_obj, target)

        if target in {"gsheet", "google_sheet", "google-sheet"}:
            values = None
            if logical.get("kind") in {"xlsx"}:
                sheets = logical.get("sheets") or {}
                first_sheet = next(iter(sheets.values()), [])
                values = first_sheet
            elif logical.get("kind") == "google_sheet":
                previews = logical.get("preview") or {}
                values = next(iter(previews.values()), [])
            elif logical.get("kind") in {"plain_text", "pdf", "docx"}:
                text = logical.get("text") or ""
                values = [[line] for line in text.splitlines()]
            if values is None:
                raise ValueError("Conversao para Google Sheet nao implementada para este formato.")
            return json.dumps({"range": "A1", "values": values}, ensure_ascii=False).encode("utf-8"), GOOGLE_NATIVE_TARGETS[target], _guess_filename_for_conversion(file_obj, target)

        if target in {"gslides", "google_slides", "google-slides"}:
            slides_payload = None
            if logical.get("kind") == "pptx":
                slides_payload = [{"title": f"Slide {slide['index']}", "body": "\n".join(slide.get("text", []))} for slide in logical.get("slides", [])]
            elif logical.get("kind") in {"plain_text", "docx", "pdf"}:
                slides_payload = [{"title": stem, "body": logical.get("text") or ""} for stem in [Path(file_obj.get("name") or "Slide").stem]]
            if slides_payload is None:
                raise ValueError("Conversao para Google Slides nao implementada para este formato.")
            return json.dumps({"slides": slides_payload}, ensure_ascii=False).encode("utf-8"), GOOGLE_NATIVE_TARGETS[target], _guess_filename_for_conversion(file_obj, target)

    if source_mime.startswith("application/vnd.google-apps"):
        export_mime = _pick_export_mime(file_obj, target_format=target, mime_type=mime_type)
        return _export_google_native_file(clients, file_obj, export_mime), export_mime, _guess_filename_for_conversion(file_obj, target, export_mime)

    binary = _download_bytes(clients, file_obj["id"])
    if not target and not mime_type:
        return binary, file_obj.get("mimeType") or _guess_mime(file_obj.get("name", "")), file_obj.get("name") or "item"

    logical = _read_file_auto(clients, file_obj)
    filename = _guess_filename_for_conversion(file_obj, target, mime_type)
    binary, resolved_mime = _serialize_content_for_filename(filename, logical.get("sheets") or logical.get("slides") or logical.get("text") or logical, mime_type)
    return binary, resolved_mime, filename


def _write_binary_file(clients: GoogleClients, file_obj: dict[str, Any], content: bytes, mime_type: str) -> dict[str, Any]:
    media = MediaInMemoryUpload(content, mimetype=mime_type, resumable=False)
    updated = clients.drive.files().update(
        fileId=file_obj["id"],
        media_body=media,
        fields="id,name,mimeType,parents,webViewLink,modifiedTime,size",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "file": _metadata(updated)}


def _write_file_auto(clients: GoogleClients, file_obj: dict[str, Any], content: Any, mode: str = "replace") -> dict[str, Any]:
    handler = _detect_handler(file_obj)
    if handler == "google_doc":
        _clear_and_write_doc(clients, file_obj["id"], str(content))
        return {"ok": True, "file": _metadata(file_obj), "mode": mode, "handler": handler}
    if handler == "google_sheet":
        if not isinstance(content, dict) or "range" not in content or "values" not in content:
            raise ValueError("Para Google Sheets, informe content como {'range': 'A1:C3', 'values': [[...]]}.")
        _sheet_write(clients, file_obj["id"], content["range"], content["values"])
        return {"ok": True, "file": _metadata(file_obj), "mode": mode, "handler": handler}
    if handler == "google_slides":
        _write_google_slides(clients, file_obj["id"], content)
        return {"ok": True, "file": _metadata(file_obj), "mode": mode, "handler": handler, "summary": _read_google_slides_summary(clients, file_obj["id"])}

    binary, mime_type = _serialize_content_for_filename(file_obj["name"], content, file_obj.get("mimeType", ""))
    result = _write_binary_file(clients, file_obj, binary, mime_type)
    result.update({"mode": mode, "handler": handler})
    return result


@mcp.tool(name="gdrive_health", title="Saude do conector Google Drive", description="Valida autenticacao, acesso a pasta-raiz e configuracao basica do conector.")
def gdrive_health() -> dict[str, Any]:
    clients = _clients()
    root = _root_metadata(clients)
    return {"ok": True, "root": _metadata(root)}


@mcp.tool(name="list_items", title="Listar itens", description="Lista arquivos e subpastas de um caminho relativo dentro da pasta-raiz do Google Drive.")
def list_items(path: str = "") -> dict[str, Any]:
    clients = _clients()
    folder = _resolve_path(clients, path, final_folder=False)
    if folder["mimeType"] != FOLDER_MIME:
        raise ValueError("O caminho informado nao e uma pasta.")
    items = [_metadata(x) for x in _list_children(clients, folder["id"])]
    return {"ok": True, "path": _normalize_relpath(path), "items": items}


@mcp.tool(name="create_folder", title="Criar pasta", description="Cria uma subpasta dentro de um caminho relativo da pasta-raiz do Google Drive.")
def create_folder(parent_path: str, name: str) -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    if parent["mimeType"] != FOLDER_MIME:
        raise ValueError("parent_path precisa apontar para uma pasta.")
    name = _slug(name)
    existing = _child_by_name(clients, parent["id"], name)
    if existing and existing["mimeType"] == FOLDER_MIME:
        return {"ok": True, "created": False, "folder": _metadata(existing)}
    folder = clients.drive.files().create(
        body={"name": name, "mimeType": FOLDER_MIME, "parents": [parent["id"]]},
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "created": True, "folder": _metadata(folder)}


@mcp.tool(name="create_doc", title="Criar Google Doc", description="Cria um Google Doc nativo dentro de uma pasta da subarvore permitida e, opcionalmente, grava conteudo inicial.")
def create_doc(parent_path: str, title: str, content: str = "") -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    doc = clients.drive.files().create(
        body={"name": _slug(title), "mimeType": GOOGLE_DOC_MIME, "parents": [parent["id"]]},
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    if content:
        _clear_and_write_doc(clients, doc["id"], content)
    return {"ok": True, "doc": _metadata(doc)}


@mcp.tool(name="read_doc", title="Ler Google Doc", description="Le o texto integral de um Google Doc nativo dentro da pasta-raiz permitida.")
def read_doc(file_id: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    if file_obj["mimeType"] != GOOGLE_DOC_MIME:
        raise ValueError("O item informado nao e um Google Doc.")
    return {"ok": True, "file": _metadata(file_obj), "text": _doc_text(clients, file_id)}


@mcp.tool(name="write_doc", title="Sobrescrever Google Doc", description="Sobrescreve integralmente o conteudo de um Google Doc nativo dentro da pasta-raiz permitida.")
def write_doc(file_id: str, content: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    if file_obj["mimeType"] != GOOGLE_DOC_MIME:
        raise ValueError("O item informado nao e um Google Doc.")
    _clear_and_write_doc(clients, file_id, content)
    return {"ok": True, "file": _metadata(file_obj)}


@mcp.tool(name="create_sheet", title="Criar Google Sheet", description="Cria uma planilha Google Sheets dentro de uma pasta da subarvore permitida.")
def create_sheet(parent_path: str, title: str, sheet_name: str = "Pagina1") -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    spreadsheet = clients.sheets.spreadsheets().create(
        body={"properties": {"title": _slug(title)}, "sheets": [{"properties": {"title": _slug(sheet_name)}}]}
    ).execute()
    file_id = spreadsheet["spreadsheetId"]
    current = _get_file(clients, file_id, fields="id,name,mimeType,parents,webViewLink")
    prev_parents = ",".join(current.get("parents", []))
    moved = clients.drive.files().update(
        fileId=file_id,
        addParents=parent["id"],
        removeParents=prev_parents,
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "sheet": _metadata(moved)}


@mcp.tool(name="read_sheet_range", title="Ler intervalo de planilha", description="Le um intervalo A1 de uma Google Sheet dentro da pasta-raiz permitida.")
def read_sheet_range(file_id: str, a1_range: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    if file_obj["mimeType"] != GOOGLE_SHEET_MIME:
        raise ValueError("O item informado nao e uma Google Sheet.")
    values = _sheet_values(clients, file_id, a1_range)
    return {"ok": True, "file": _metadata(file_obj), "range": a1_range, "values": values}


@mcp.tool(name="write_sheet_range", title="Gravar intervalo de planilha", description="Escreve valores em um intervalo A1 de uma Google Sheet dentro da pasta-raiz permitida.")
def write_sheet_range(file_id: str, a1_range: str, values: list[list[Any]]) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    if file_obj["mimeType"] != GOOGLE_SHEET_MIME:
        raise ValueError("O item informado nao e uma Google Sheet.")
    _sheet_write(clients, file_id, a1_range, values)
    return {"ok": True, "file": _metadata(file_obj), "range": a1_range}


@mcp.tool(name="upload_text_file", title="Salvar arquivo de texto simples", description="Cria ou substitui um arquivo de texto simples dentro de uma pasta da subarvore permitida.")
def upload_text_file(parent_path: str, filename: str, content: str, mime_type: str = DEFAULT_TEXT_MIME, replace_if_exists: bool = True) -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    return _upload_or_replace_bytes(clients, parent, filename, content.encode("utf-8"), mime_type, replace_if_exists)


@mcp.tool(name="upload_base64_file", title="Salvar arquivo binario em base64", description="Cria ou substitui um arquivo binario pequeno a partir de conteudo base64, dentro da subarvore permitida.")
def upload_base64_file(parent_path: str, filename: str, base64_content: str, mime_type: str = "application/octet-stream", replace_if_exists: bool = True) -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    return _upload_or_replace_bytes(clients, parent, filename, base64.b64decode(base64_content), mime_type or _guess_mime(filename), replace_if_exists)


@mcp.tool(name="download_file_metadata", title="Ler metadados de arquivo", description="Le metadados de um item dentro da subarvore permitida do Google Drive.")
def download_file_metadata(file_id: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    return {"ok": True, "file": _metadata(file_obj)}


@mcp.tool(name="get_file_metadata", title="Obter metadados de arquivo", description="Alias conceitual de alto nivel para leitura de metadados de um item na subarvore permitida.")
def get_file_metadata(file_id: str) -> dict[str, Any]:
    return download_file_metadata(file_id)


@mcp.tool(name="read_text_file", title="Ler arquivo de texto simples", description="Le o conteudo de um arquivo de texto simples dentro da pasta-raiz permitida.")
def read_text_file(file_id: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    mime_type = file_obj.get("mimeType", "")
    if mime_type.startswith("application/vnd.google-apps"):
        raise ValueError("Use as ferramentas nativas de Google Docs ou Google Sheets para este arquivo.")
    text = _download_bytes(clients, file_id).decode("utf-8", errors="replace")
    return {"ok": True, "file": _metadata(file_obj), "text": text}


@mcp.tool(name="read_file", title="Ler arquivo", description="Le um arquivo da subarvore permitida, escolhendo automaticamente o melhor handler por MIME ou extensao.")
def read_file(file_id: str, mode: str = "auto") -> dict[str, Any]:
    if mode != "auto":
        raise ValueError("Nesta fase, apenas mode='auto' esta implementado.")
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    content = _read_file_auto(clients, file_obj)
    return {"ok": True, "file": _metadata(file_obj), "handler": _detect_handler(file_obj), "content": content}


@mcp.tool(name="create_file", title="Criar arquivo", description="Cria um arquivo na pasta destino preservando o formato pedido quando isso for tecnicamente viavel.")
def create_file(parent_path: str, filename: str, content: Any = "", mime_type: str = "", format_hint: str = "", replace_if_exists: bool = False) -> dict[str, Any]:
    clients = _clients()
    parent = _resolve_path(clients, parent_path, final_folder=False)
    name = _slug(filename)
    ext = _ext(name)
    hint = (format_hint or "").lower().strip()
    chosen_mime = (mime_type or "").strip()

    if hint in {"gdoc", "google_doc", "google-doc"} or chosen_mime == GOOGLE_DOC_MIME:
        return create_doc(parent_path=parent_path, title=name, content=str(content))
    if hint in {"gsheet", "google_sheet", "google-sheet"} or chosen_mime == GOOGLE_SHEET_MIME:
        sheet = create_sheet(parent_path=parent_path, title=name)
        if content:
            file_id = sheet["sheet"]["id"]
            if isinstance(content, dict) and "range" in content and "values" in content:
                _sheet_write(clients, file_id, content["range"], content["values"])
        return sheet
    if hint in {"gslides", "google_slides", "google-slides"} or chosen_mime == GOOGLE_SLIDES_MIME:
        return _create_google_slides(clients, parent_path=parent_path, title=name, slides=content)

    binary, resolved_mime = _serialize_content_for_filename(name, content, chosen_mime)
    return _upload_or_replace_bytes(clients, parent, name, binary, resolved_mime, replace_if_exists)


@mcp.tool(name="write_file", title="Escrever arquivo", description="Sobrescreve um arquivo usando o handler mais adequado para o formato detectado.")
def write_file(file_id: str, content: Any, mode: str = "replace") -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    return _write_file_auto(clients, file_obj, content, mode)


@mcp.tool(name="convert_file", title="Converter arquivo", description="Converte um arquivo para outro formato ou para um formato Google nativo, preservando-o na subarvore permitida.")
def convert_file(file_id: str, target_format: str = "", target_mime_type: str = "", destination_parent_path: str = "", output_filename: str = "", replace_if_exists: bool = False) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    if destination_parent_path:
        parent = _resolve_path(clients, destination_parent_path, final_folder=False)
    else:
        parent_id = (file_obj.get("parents") or [ROOT_FOLDER_ID])[0]
        parent = _get_file(clients, parent_id, fields="id,name,mimeType,parents,webViewLink")

    binary, resolved_mime, suggested_name = _materialize_converted_content(
        clients,
        file_obj,
        target_format=target_format,
        mime_type=target_mime_type,
    )

    final_name = _slug(output_filename or suggested_name or file_obj.get("name") or "item")
    target = (target_format or "").lower().strip()

    if resolved_mime == GOOGLE_DOC_MIME or target in {"gdoc", "google_doc", "google-doc"}:
        created = clients.drive.files().create(
            body={"name": final_name, "mimeType": GOOGLE_DOC_MIME, "parents": [parent["id"]]},
            fields="id,name,mimeType,parents,webViewLink",
            supportsAllDrives=True,
        ).execute()
        _clear_and_write_doc(clients, created["id"], binary.decode("utf-8", errors="replace"))
        return {"ok": True, "source": _metadata(file_obj), "target": _metadata(created), "targetFormat": target or resolved_mime}

    if resolved_mime == GOOGLE_SHEET_MIME or target in {"gsheet", "google_sheet", "google-sheet"}:
        payload = json.loads(binary.decode("utf-8"))
        spreadsheet = clients.sheets.spreadsheets().create(
            body={"properties": {"title": final_name}, "sheets": [{"properties": {"title": "Pagina1"}}]}
        ).execute()
        new_id = spreadsheet["spreadsheetId"]
        current = _get_file(clients, new_id, fields="id,name,mimeType,parents,webViewLink")
        prev_parents = ",".join(current.get("parents", []))
        moved = clients.drive.files().update(
            fileId=new_id,
            addParents=parent["id"],
            removeParents=prev_parents,
            fields="id,name,mimeType,parents,webViewLink",
            supportsAllDrives=True,
        ).execute()
        _sheet_write(clients, new_id, payload.get("range", "A1"), payload.get("values", []))
        return {"ok": True, "source": _metadata(file_obj), "target": _metadata(moved), "targetFormat": target or resolved_mime}

    if resolved_mime == GOOGLE_SLIDES_MIME or target in {"gslides", "google_slides", "google-slides"}:
        payload = json.loads(binary.decode("utf-8"))
        created = clients.drive.files().create(
            body={"name": final_name, "mimeType": GOOGLE_SLIDES_MIME, "parents": [parent["id"]]},
            fields="id,name,mimeType,parents,webViewLink",
            supportsAllDrives=True,
        ).execute()
        _write_google_slides(clients, created["id"], payload.get("slides", []))
        return {"ok": True, "source": _metadata(file_obj), "target": _metadata(created), "targetFormat": target or resolved_mime, "slidesRequested": payload.get("slides", []), "summary": _read_google_slides_summary(clients, created["id"])}

    created = _upload_or_replace_bytes(clients, parent, final_name, binary, resolved_mime, replace_if_exists)
    return {"ok": True, "source": _metadata(file_obj), "target": created.get("file"), "targetFormat": target or resolved_mime}


@mcp.tool(name="move_item", title="Mover item", description="Move um arquivo ou subpasta para outra pasta da subarvore permitida.")
def move_item(file_id: str, destination_folder_path: str) -> dict[str, Any]:
    clients = _clients()
    file_obj = _ensure_inside_root(clients, file_id)
    dest = _resolve_path(clients, destination_folder_path, final_folder=False)
    prev_parents = ",".join(file_obj.get("parents", []))
    moved = clients.drive.files().update(
        fileId=file_id,
        addParents=dest["id"],
        removeParents=prev_parents,
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "file": _metadata(moved)}


@mcp.tool(name="rename_item", title="Renomear item", description="Renomeia um arquivo ou pasta dentro da subarvore permitida.")
def rename_item(file_id: str, new_name: str) -> dict[str, Any]:
    clients = _clients()
    _ensure_inside_root(clients, file_id)
    updated = clients.drive.files().update(
        fileId=file_id,
        body={"name": _slug(new_name)},
        fields="id,name,mimeType,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "file": _metadata(updated)}


@mcp.tool(name="trash_item", title="Enviar item para a lixeira", description="Envia um arquivo ou pasta para a lixeira dentro da subarvore permitida.")
def trash_item(file_id: str) -> dict[str, Any]:
    clients = _clients()
    _ensure_inside_root(clients, file_id)
    trashed = clients.drive.files().update(
        fileId=file_id,
        body={"trashed": True},
        fields="id,name,mimeType,trashed,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "file": _metadata(trashed)}


@mcp.tool(name="restore_item", title="Restaurar item da lixeira", description="Restaura um arquivo ou pasta da lixeira, desde que o item pertença a subarvore permitida.")
def restore_item(file_id: str) -> dict[str, Any]:
    clients = _clients()
    _ensure_inside_root(clients, file_id)
    restored = clients.drive.files().update(
        fileId=file_id,
        body={"trashed": False},
        fields="id,name,mimeType,trashed,parents,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return {"ok": True, "file": _metadata(restored)}


@mcp.tool(name="search_items", title="Pesquisar itens", description="Pesquisa itens por nome dentro da subarvore permitida do Google Drive. Opcionalmente filtra por mime type e pode buscar recursivamente.")
def search_items(query: str, path: str = "", mime_type: str = "", recursive: bool = True, max_results: int = 100) -> dict[str, Any]:
    clients = _clients()
    folder = _resolve_path(clients, path, final_folder=False)
    query = (query or "").strip()
    if not query:
        raise ValueError("query e obrigatoria.")

    root_folder_id = folder["id"]
    folder_ids = [root_folder_id]
    seen = {root_folder_id}

    if recursive:
        idx = 0
        while idx < len(folder_ids):
            current_id = folder_ids[idx]
            idx += 1
            for child in _list_children(clients, current_id):
                if child.get("mimeType") == FOLDER_MIME and child["id"] not in seen:
                    seen.add(child["id"])
                    folder_ids.append(child["id"])

    tokens = [token.replace("'", "\'") for token in query.split() if token.strip()]
    results: list[dict[str, Any]] = []
    matched_ids: set[str] = set()

    for current_id in folder_ids:
        clauses = [f"'{current_id}' in parents", "trashed = false"]
        for token in tokens:
            clauses.append(f"name contains '{token}'")
        if mime_type:
            safe_mime = mime_type.replace("'", "\'")
            clauses.append(f"mimeType = '{safe_mime}'")
        resp = clients.drive.files().list(
            q=" and ".join(clauses),
            fields="files(id,name,mimeType,parents,modifiedTime,size,webViewLink)",
            pageSize=min(max_results, 1000),
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
        ).execute()
        for item in resp.get("files", []):
            if item["id"] in matched_ids:
                continue
            matched_ids.add(item["id"])
            results.append(_metadata(item))
            if len(results) >= max_results:
                return {"ok": True, "items": results, "recursive": recursive, "searchedFolders": len(folder_ids), "truncated": True}

    return {"ok": True, "items": results, "recursive": recursive, "searchedFolders": len(folder_ids), "truncated": False}



if __name__ == "__main__":
    mcp.run(transport="streamable-http")
